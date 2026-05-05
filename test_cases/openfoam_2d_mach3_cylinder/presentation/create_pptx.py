#!/usr/bin/env python3
"""
Build PowerPoint presentation for Mach 3 Supersonic Cylinder AMR results.
Fahim Shahriyar — University of North Dakota
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = "/home/fahim.shahriyar/dsmc/openfoam_mach3_cylinder/presentation"
IMG  = os.path.join(BASE, "images")

# Colors
DARK_BLUE   = RGBColor(0x00, 0x2B, 0x5C)  # UND dark blue
MEDIUM_BLUE = RGBColor(0x00, 0x4E, 0x8C)
LIGHT_BLUE  = RGBColor(0x00, 0x7B, 0xBB)
ACCENT_GREEN = RGBColor(0x00, 0x82, 0x50)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
TABLE_HEADER = RGBColor(0x00, 0x2B, 0x5C)
TABLE_ALT    = RGBColor(0xF0, 0xF4, 0xF8)
BOX_BG       = RGBColor(0xF5, 0xF7, 0xFA)
ARROW_COLOR  = RGBColor(0x00, 0x4E, 0x8C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========================================================================
# Helper Functions
# ========================================================================
def add_background(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a dark blue title bar at the top."""
    # Title bar background
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                  prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title text
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
        p2.alignment = PP_ALIGN.LEFT

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.2),
                                   prs.slide_width, Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()

def add_footer(slide, slide_num, total=17):
    """Add footer with slide number and author."""
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(6), Inches(0.35))
    tf = ft.text_frame
    p = tf.paragraphs[0]
    p.text = "Fahim Shahriyar | University of North Dakota | Dr. Chonglin Zhang"
    p.font.size = Pt(10)
    p.font.color.rgb = MED_GRAY

    sn = slide.shapes.add_textbox(Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.35))
    tf2 = sn.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"{slide_num} / {total}"
    p2.font.size = Pt(10)
    p2.font.color.rgb = MED_GRAY
    p2.alignment = PP_ALIGN.RIGHT

def add_bullet_text(slide, left, top, width, height, items, font_size=16):
    """Add bulleted text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
        p.level = 0

def add_info_box(slide, left, top, width, height, title, items, title_color=DARK_BLUE):
    """Add a styled information box with title and bullet items."""
    # Background box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = BOX_BG
    box.line.color.rgb = RGBColor(0xCC, 0xD5, 0xE0)
    box.line.width = Pt(1)

    # Title
    ttl = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1),
                                    width - Inches(0.4), Inches(0.4))
    tf = ttl.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = title_color

    # Items
    itb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.5),
                                    width - Inches(0.5), height - Inches(0.6))
    tf2 = itb.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        # Handle bold parts (wrapped in **)
        if "**" in item:
            parts = item.split("**")
            for j, part in enumerate(parts):
                run = p.add_run()
                run.text = part
                run.font.size = Pt(14)
                run.font.color.rgb = DARK_GRAY
                if j % 2 == 1:
                    run.font.bold = True
        else:
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

def add_table(slide, left, top, width, rows_data, col_widths=None):
    """Add a styled table."""
    nrows = len(rows_data)
    ncols = len(rows_data[0])
    table_shape = slide.shapes.add_table(nrows, ncols, left, top, width,
                                          Inches(0.4 * nrows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape

TOTAL_SLIDES = 17

# ========================================================================
# SLIDE 1: Title Slide
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_background(slide)

# Big blue block
block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                prs.slide_width, Inches(4.5))
block.fill.solid()
block.fill.fore_color.rgb = DARK_BLUE
block.line.fill.background()

# Title
ttl = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(11), Inches(1.5))
tf = ttl.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Adaptive Mesh Refinement for Supersonic Flow"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Mach 3 Cylinder at 50 km Altitude"
p2.font.size = Pt(28)
p2.font.color.rgb = RGBColor(0x88, 0xBB, 0xEE)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(12)

p3 = tf.add_paragraph()
p3.text = "OpenFOAM v2406 | Gradient-Based Adaptation | Unstructured Gmsh Mesh"
p3.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0xAA, 0xCC, 0xDD)
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(20)

# Author info below blue block
auth = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11), Inches(2))
tf2 = auth.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Fahim Shahriyar"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER

p2 = tf2.add_paragraph()
p2.text = "M.S. Mechanical Engineering"
p2.font.size = Pt(18)
p2.font.color.rgb = MED_GRAY
p2.alignment = PP_ALIGN.CENTER

p3 = tf2.add_paragraph()
p3.text = "Advisor: Dr. Chonglin Zhang"
p3.font.size = Pt(18)
p3.font.color.rgb = MED_GRAY
p3.alignment = PP_ALIGN.CENTER

p4 = tf2.add_paragraph()
p4.text = "Department of Mechanical Engineering, University of North Dakota"
p4.font.size = Pt(16)
p4.font.color.rgb = MED_GRAY
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(8)

# ========================================================================
# SLIDE 2: Outline
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Outline")
add_footer(slide, 2, TOTAL_SLIDES)

outline_items = [
    ("1.", "Problem Statement & Motivation"),
    ("2.", "Flow Conditions (50 km Altitude)"),
    ("3.", "Computational Setup & Boundary Conditions"),
    ("4.", "Solver & Numerical Schemes"),
    ("5.", "AMR Pipeline Architecture"),
    ("6.", "Gradient-Based Sizing Formula"),
    ("7.", "Mesh Progression (Uniform through AMR 3)"),
    ("8.", "Pressure Contours"),
    ("9.", "Temperature Contours"),
    ("10.", "Velocity Contours"),
    ("11.", "Mesh & Results Comparison"),
    ("12.", "Key Findings & Conclusions"),
]

for i, (num, text) in enumerate(outline_items):
    y = Inches(1.6) + Inches(i * 0.42)
    nb = slide.shapes.add_textbox(Inches(2.5), y, Inches(1), Inches(0.4))
    tf = nb.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.RIGHT

    tb = slide.shapes.add_textbox(Inches(3.5), y, Inches(7), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY

# ========================================================================
# SLIDE 3: Problem Statement & Motivation
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Problem Statement & Motivation")
add_footer(slide, 3, TOTAL_SLIDES)

add_info_box(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(2.5),
    "Problem",
    [
        "Supersonic flows produce strong shock waves with steep gradients",
        "Uniform meshes waste cells in smooth regions far from shocks",
        "Fine uniform meshes are computationally expensive",
        "Need automated refinement that focuses cells where needed",
    ])

add_info_box(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(2.5),
    "Objective",
    [
        "Develop automated AMR pipeline for unstructured meshes",
        "Use pressure gradient as adaptation criterion",
        "Demonstrate on Mach 3 cylinder at 50 km altitude",
        "Pipeline: Gmsh + OpenFOAM + ParaView (fully automated)",
    ])

add_info_box(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.3),
    "Why This Matters",
    [
        "Supersonic/hypersonic vehicle design requires accurate shock resolution without excessive computational cost",
        "Gradient-based adaptation automatically concentrates cells at bow shocks, wake regions, and stagnation points",
        "The pipeline is solver-agnostic: works with OpenFOAM (CFD) and COMET (DSMC) through a unified input file",
        "Single configuration file (amr_pipeline.input) controls the entire process; no script editing required",
    ])

# ========================================================================
# SLIDE 4: Flow Conditions
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Flow Conditions", "1976 US Standard Atmosphere at 50 km Altitude")
add_footer(slide, 4, TOTAL_SLIDES)

# Left: Freestream conditions table
rows = [
    ["Parameter", "Value", "Unit"],
    ["Altitude", "50", "km"],
    ["Mach Number", "3.0", "-"],
    ["Freestream Velocity", "990", "m/s"],
    ["Freestream Pressure", "79.78", "Pa"],
    ["Freestream Temperature", "270.65", "K"],
    ["Freestream Density", "1.027 x 10^-3", "kg/m^3"],
    ["Speed of Sound", "330", "m/s"],
    ["Wall Temperature", "800", "K"],
    ["Turbulence Model", "Laminar", "-"],
]
add_table(slide, Inches(0.5), Inches(1.6), Inches(6.5), rows,
          [Inches(2.5), Inches(2.5), Inches(1.5)])

# Right: key notes
add_info_box(slide, Inches(7.3), Inches(1.6), Inches(5.5), Inches(3.0),
    "Notes",
    [
        "Density computed: rho = p / (R * T) = 79.78 / (287.05 * 270.65)",
        "Speed of sound: a = sqrt(gamma * R * T) = sqrt(1.4 * 287.05 * 270.65) = 330 m/s",
        "Mach number: M = V/a = 990/330 = 3.0",
        "Reference: NASA-TM-X-74335 (1976 US Standard Atmosphere)",
    ])

add_info_box(slide, Inches(7.3), Inches(4.9), Inches(5.5), Inches(1.8),
    "Domain Geometry",
    [
        "Rectangular domain: 1.0 m x 0.8 m",
        "Cylinder diameter: 0.2 m, centered at (0.5, 0.4)",
        "Quasi-2D: single cell in z-direction (0.01 m extrude)",
        "Unstructured triangular mesh generated by Gmsh",
    ])

# ========================================================================
# SLIDE 5: Computational Setup & BCs
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Computational Setup & Boundary Conditions")
add_footer(slide, 5, TOTAL_SLIDES)

# BC table
bc_rows = [
    ["Boundary", "Velocity (U)", "Pressure (p)", "Temperature (T)"],
    ["Inlet (left)", "fixedValue: (990, 0, 0) m/s", "fixedValue: 79.78 Pa", "fixedValue: 270.65 K"],
    ["Outlet (bottom,\nright, top)", "zeroGradient", "zeroGradient", "zeroGradient"],
    ["Wall (cylinder)", "noSlip", "zeroGradient", "fixedValue: 800 K"],
    ["Front & Back", "empty", "empty", "empty"],
]
add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3), bc_rows,
          [Inches(2.0), Inches(3.4), Inches(3.4), Inches(3.5)])

add_info_box(slide, Inches(0.5), Inches(4.0), Inches(6), Inches(3.0),
    "Initial Conditions",
    [
        "Internal field initialized to freestream values",
        "U = (990, 0, 0) m/s everywhere",
        "p = 79.78 Pa everywhere",
        "T = 270.65 K everywhere",
        "Simulation starts from uniform initial state",
    ])

add_info_box(slide, Inches(6.8), Inches(4.0), Inches(6), Inches(3.0),
    "Gas Properties (thermophysicalProperties)",
    [
        "Equation of state: Perfect gas (hePsiThermo)",
        "Molecular weight: 28.9 g/mol (air)",
        "Specific heat: Cp = 1005 J/(kg K)",
        "Transport: Sutherland law (As=1.458e-6, Ts=110.4 K)",
        "gamma = 1.4 (ratio of specific heats)",
    ])

# ========================================================================
# SLIDE 6: Solver & Numerical Schemes
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Solver & Numerical Schemes")
add_footer(slide, 6, TOTAL_SLIDES)

add_info_box(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.8),
    "Solver: rhoCentralFoam",
    [
        "Density-based compressible solver",
        "Kurganov-Tadmor central scheme for flux evaluation",
        "Semi-discrete, non-staggered formulation",
        "Suitable for high-speed flows with strong shocks",
        "No Riemann solver needed (central differencing)",
    ])

add_info_box(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(2.8),
    "Reconstruction & Limiters",
    [
        "Density (rho): Minmod limiter",
        "Velocity (U): MinmodV limiter (vector)",
        "Temperature (T): Minmod limiter",
        "Minmod is most diffusive but most stable for strong shocks",
        "Prevents spurious oscillations near discontinuities",
    ])

# Time stepping
rows2 = [
    ["Parameter", "Value"],
    ["Time integration", "Euler (1st order)"],
    ["Initial time step", "1e-8 s"],
    ["Max Courant number", "0.3"],
    ["Max time step", "1e-5 s"],
    ["End time", "0.001 s"],
    ["Write interval", "0.0002 s"],
    ["Adaptive time stepping", "Yes"],
]
add_table(slide, Inches(0.5), Inches(4.7), Inches(5.8), rows2,
          [Inches(3.0), Inches(2.8)])

add_info_box(slide, Inches(6.8), Inches(4.7), Inches(6), Inches(2.3),
    "Why These Choices?",
    [
        "rhoCentralFoam: Only OpenFOAM solver that handles Mach 3 on unstructured mesh",
        "sonicFoam crashes at ANY supersonic speed on triangular mesh",
        "Minmod: Most stable limiter; vanLeer/vanAlbada produce oscillations at Mach 3",
        "maxCo = 0.3: Required for stability with strong bow shock",
        "Euler time integration: Sufficient for steady-state convergence",
    ])

# ========================================================================
# SLIDE 7: AMR Pipeline Architecture (Flowchart)
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "AMR Pipeline Architecture", "Automated Workflow: all_run.sh")
add_footer(slide, 7, TOTAL_SLIDES)

# Draw flowchart boxes and arrows
box_w = Inches(2.0)
box_h = Inches(0.8)
arr_len = Inches(0.4)

steps = [
    ("Gmsh", "Generate mesh\n(.geo -> .msh)", MEDIUM_BLUE),
    ("gmshToFoam", "Convert mesh\n(.msh -> polyMesh)", MEDIUM_BLUE),
    ("Boundary Fix", "Set patch types\n(empty, wall)", MEDIUM_BLUE),
    ("rhoCentralFoam", "Run CFD solver\n(Kurganov scheme)", RGBColor(0xB8, 0x40, 0x00)),
    ("foamToVTK", "Export results\n(-> .vtu files)", MEDIUM_BLUE),
    ("pvpython", "Compute grad(p)\nSizing formula", ACCENT_GREEN),
    ("Export .pos", "Background sizing\nfield for Gmsh", ACCENT_GREEN),
]

# Two rows: top row (steps 1-4), bottom row (steps 5-7 + loop back)
# Top row
start_x = Inches(0.4)
y_top = Inches(1.8)
y_bot = Inches(4.2)

positions = []
for i, (title, desc, color) in enumerate(steps[:4]):
    x = start_x + i * (box_w + arr_len + Inches(0.3))
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_top, box_w, box_h + Inches(0.3))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    p2.alignment = PP_ALIGN.CENTER
    positions.append((x, y_top))

    # Arrow to next
    if i < 3:
        ax = x + box_w + Inches(0.05)
        ay = y_top + (box_h + Inches(0.3)) / 2 - Inches(0.1)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, Inches(0.5), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ARROW_COLOR
        arrow.line.fill.background()

# Down arrow from step 4
ax4 = positions[3][0] + box_w / 2 - Inches(0.12)
ay4 = y_top + box_h + Inches(0.3) + Inches(0.05)
darr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, ax4, ay4, Inches(0.25), Inches(0.5))
darr.fill.solid()
darr.fill.fore_color.rgb = ARROW_COLOR
darr.line.fill.background()

# Bottom row (right to left: steps 5, 6, 7)
for i, (title, desc, color) in enumerate(steps[4:]):
    x = positions[3][0] - i * (box_w + arr_len + Inches(0.3))
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_bot, box_w, box_h + Inches(0.3))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    p2.alignment = PP_ALIGN.CENTER

    # Left arrow between bottom boxes
    if i > 0:
        prev_x = positions[3][0] - (i - 1) * (box_w + arr_len + Inches(0.3))
        ax = x + box_w + Inches(0.05)
        ay = y_bot + (box_h + Inches(0.3)) / 2 - Inches(0.1)
        arrow = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, ax, ay, Inches(0.5), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ARROW_COLOR
        arrow.line.fill.background()

# Loop-back arrow label
loop_x = positions[3][0] - 2 * (box_w + arr_len + Inches(0.3))
loop_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   loop_x - Inches(1.5), y_bot + Inches(0.1),
                                   Inches(1.3), Inches(0.9))
loop_box.fill.solid()
loop_box.fill.fore_color.rgb = RGBColor(0xCC, 0x33, 0x33)
loop_box.line.fill.background()
tf = loop_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Loop Back\nto Step 1"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

larr = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW,
    loop_x - Inches(0.15), y_bot + (box_h + Inches(0.3)) / 2 - Inches(0.1),
    Inches(0.4), Inches(0.25))
larr.fill.solid()
larr.fill.fore_color.rgb = RGBColor(0xCC, 0x33, 0x33)
larr.line.fill.background()

# Config file callout
cfg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.4), Inches(5.8), Inches(12.3), Inches(1.0))
cfg.fill.solid()
cfg.fill.fore_color.rgb = BOX_BG
cfg.line.color.rgb = ACCENT_GREEN
cfg.line.width = Pt(2)
tf = cfg.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Configuration: "
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = ACCENT_GREEN
run = p.add_run()
run.text = ("amr_pipeline.input controls everything: loops=3, extraction_mode=gradient, "
            "gradient_field=p, sizing_min=0.002, sizing_max=0.015, sizing_scale=100.0, "
            "solver=rhoCentralFoam. No script editing required.")
run.font.size = Pt(13)
run.font.color.rgb = DARK_GRAY

# ========================================================================
# SLIDE 8: Sizing Formula
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Gradient-Based Mesh Sizing Formula")
add_footer(slide, 8, TOTAL_SLIDES)

# Formula box
fbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(1.5), Inches(1.8), Inches(10), Inches(1.8))
fbox.fill.solid()
fbox.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFF)
fbox.line.color.rgb = MEDIUM_BLUE
fbox.line.width = Pt(2)

tf = fbox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "h(x) = h_min + (h_max - h_min) / (1 + alpha * |grad(p)| / |grad(p)|_max)"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(20)

p2 = tf.add_paragraph()
p2.text = "Frey & Alauzet (2005), Loseille & Alauzet (2011)"
p2.font.size = Pt(14)
p2.font.italic = True
p2.font.color.rgb = MED_GRAY
p2.alignment = PP_ALIGN.CENTER

# Parameter table
param_rows = [
    ["Symbol", "Parameter", "Value", "Description"],
    ["h_min", "Minimum cell size", "0.002 m", "Smallest allowed element in refined regions"],
    ["h_max", "Maximum cell size", "0.015 m", "Largest allowed element in smooth regions"],
    ["alpha", "Scaling factor", "100.0", "Controls sensitivity to gradient magnitude"],
    ["grad(p)", "Gradient field", "Pressure", "Chosen adaptation criterion"],
]
add_table(slide, Inches(0.5), Inches(4.0), Inches(12.3), param_rows,
          [Inches(1.5), Inches(2.5), Inches(2.0), Inches(6.3)])

add_info_box(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2),
    "How It Works",
    [
        "High pressure gradient (shock region): h approaches h_min = 0.002 m (fine cells concentrate at bow shock, wake, stagnation)",
        "Low pressure gradient (smooth flow): h approaches h_max = 0.015 m (coarse cells save computation in freestream)",
    ])

# ========================================================================
# SLIDE 9: Mesh Progression (2x2 grid)
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Mesh Progression", "Uniform Mesh through 3 AMR Iterations")
add_footer(slide, 9, TOTAL_SLIDES)

mesh_labels = [
    ("Uniform (4,804 nodes)", "mesh_uniform.png"),
    ("AMR 1 (71,904 nodes)", "mesh_amr1.png"),
    ("AMR 2 (14,888 nodes)", "mesh_amr2.png"),
    ("AMR 3 (17,500 nodes)", "mesh_amr3.png"),
]

img_w = Inches(5.8)
img_h = Inches(2.6)
x_positions = [Inches(0.4), Inches(6.8)]
y_positions = [Inches(1.5), Inches(4.3)]

for i, (label, fname) in enumerate(mesh_labels):
    x = x_positions[i % 2]
    y = y_positions[i // 2]

    fpath = os.path.join(IMG, fname)
    if os.path.exists(fpath):
        slide.shapes.add_picture(fpath, x, y + Inches(0.3), img_w, img_h)

    lbl = slide.shapes.add_textbox(x, y, img_w, Inches(0.35))
    tf = lbl.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

# ========================================================================
# SLIDE 10: Pressure Contours (2x2)
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Pressure Field Comparison", "Consistent scale: 0 - 960 Pa")
add_footer(slide, 10, TOTAL_SLIDES)

p_labels = [
    ("Uniform Mesh", "p_uniform.png"),
    ("AMR 1", "p_amr1.png"),
    ("AMR 2", "p_amr2.png"),
    ("AMR 3", "p_amr3.png"),
]
for i, (label, fname) in enumerate(p_labels):
    x = x_positions[i % 2]
    y = y_positions[i // 2]
    fpath = os.path.join(IMG, fname)
    if os.path.exists(fpath):
        slide.shapes.add_picture(fpath, x, y + Inches(0.3), img_w, img_h)
    lbl = slide.shapes.add_textbox(x, y, img_w, Inches(0.35))
    tf = lbl.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

# ========================================================================
# SLIDE 11: Temperature Contours (2x2)
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Temperature Field Comparison", "Consistent scale: 200 - 900 K")
add_footer(slide, 11, TOTAL_SLIDES)

t_labels = [
    ("Uniform Mesh", "T_uniform.png"),
    ("AMR 1", "T_amr1.png"),
    ("AMR 2", "T_amr2.png"),
    ("AMR 3", "T_amr3.png"),
]
for i, (label, fname) in enumerate(t_labels):
    x = x_positions[i % 2]
    y = y_positions[i // 2]
    fpath = os.path.join(IMG, fname)
    if os.path.exists(fpath):
        slide.shapes.add_picture(fpath, x, y + Inches(0.3), img_w, img_h)
    lbl = slide.shapes.add_textbox(x, y, img_w, Inches(0.35))
    tf = lbl.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

# ========================================================================
# SLIDE 12: Velocity Contours (2x2)
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Velocity Magnitude Comparison", "Consistent scale: 0 - 1100 m/s")
add_footer(slide, 12, TOTAL_SLIDES)

u_labels = [
    ("Uniform Mesh", "U_uniform.png"),
    ("AMR 1", "U_amr1.png"),
    ("AMR 2", "U_amr2.png"),
    ("AMR 3", "U_amr3.png"),
]
for i, (label, fname) in enumerate(u_labels):
    x = x_positions[i % 2]
    y = y_positions[i // 2]
    fpath = os.path.join(IMG, fname)
    if os.path.exists(fpath):
        slide.shapes.add_picture(fpath, x, y + Inches(0.3), img_w, img_h)
    lbl = slide.shapes.add_textbox(x, y, img_w, Inches(0.35))
    tf = lbl.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

# ========================================================================
# SLIDE 13: Detailed Pressure Comparison (individual larger images)
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Pressure: Uniform vs AMR 3", "Shock resolution improvement with gradient-based adaptation")
add_footer(slide, 13, TOTAL_SLIDES)

big_w = Inches(5.8)
big_h = Inches(4.5)
fpath1 = os.path.join(IMG, "p_uniform.png")
fpath2 = os.path.join(IMG, "p_amr3.png")
if os.path.exists(fpath1):
    slide.shapes.add_picture(fpath1, Inches(0.4), Inches(1.8), big_w, big_h)
if os.path.exists(fpath2):
    slide.shapes.add_picture(fpath2, Inches(6.8), Inches(1.8), big_w, big_h)

for x, label in [(Inches(0.4), "Uniform (4,804 nodes)"), (Inches(6.8), "AMR 3 (17,500 nodes)")]:
    lbl = slide.shapes.add_textbox(x, Inches(1.45), big_w, Inches(0.35))
    tf = lbl.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

note = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.4))
tf = note.text_frame
p = tf.paragraphs[0]
p.text = "AMR 3 resolves the bow shock structure with sharper gradients using 3.6x more nodes concentrated at high-gradient regions"
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = MED_GRAY
p.alignment = PP_ALIGN.CENTER

# ========================================================================
# SLIDE 14: Detailed Temperature Comparison
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Temperature: Uniform vs AMR 3", "Thermal structure near cylinder wall and wake")
add_footer(slide, 14, TOTAL_SLIDES)

fpath1 = os.path.join(IMG, "T_uniform.png")
fpath2 = os.path.join(IMG, "T_amr3.png")
if os.path.exists(fpath1):
    slide.shapes.add_picture(fpath1, Inches(0.4), Inches(1.8), big_w, big_h)
if os.path.exists(fpath2):
    slide.shapes.add_picture(fpath2, Inches(6.8), Inches(1.8), big_w, big_h)

for x, label in [(Inches(0.4), "Uniform (4,804 nodes)"), (Inches(6.8), "AMR 3 (17,500 nodes)")]:
    lbl = slide.shapes.add_textbox(x, Inches(1.45), big_w, Inches(0.35))
    tf = lbl.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

note = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.4))
tf = note.text_frame
p = tf.paragraphs[0]
p.text = "Temperature peaks at stagnation point (~900 K) due to adiabatic compression; AMR captures thermal layer thickness more accurately"
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = MED_GRAY
p.alignment = PP_ALIGN.CENTER

# ========================================================================
# SLIDE 15: Mesh Statistics Table
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Mesh & Results Summary")
add_footer(slide, 15, TOTAL_SLIDES)

stats_rows = [
    ["Mesh", "Nodes", "Cells", "Mesh File Size", "p Range (Pa)", "T Range (K)"],
    ["Uniform", "4,804", "4,592", "652 KB", "9.3 - 953.3", "256.5 - 883.0"],
    ["AMR 1", "71,904", "71,275", "11 MB", "11.8 - 973.7", "233.0 - 864.0"],
    ["AMR 2", "14,888", "14,528", "2.2 MB", "12.4 - 961.5", "237.6 - 912.7"],
    ["AMR 3", "17,500", "17,141", "2.6 MB", "12.6 - 970.9", "241.0 - 907.9"],
]
add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3), stats_rows,
          [Inches(1.8), Inches(1.8), Inches(1.8), Inches(2.0), Inches(2.4), Inches(2.5)])

add_info_box(slide, Inches(0.5), Inches(3.8), Inches(6), Inches(3.0),
    "Observations",
    [
        "AMR 1 over-refines (71k nodes) because initial gradients are everywhere",
        "AMR 2 corrects by re-solving on AMR 1 mesh; now gradients are well-resolved",
        "AMR 3 adds refinement only where residual gradients exist (17.5k nodes)",
        "Final mesh (AMR 3) is 3.6x larger than uniform but 4.1x smaller than AMR 1",
        "Pressure/temperature ranges converge across iterations",
    ])

add_info_box(slide, Inches(6.8), Inches(3.8), Inches(6), Inches(3.0),
    "AMR Pipeline Parameters",
    [
        "Adaptation criterion: Pressure gradient |grad(p)|",
        "h_min = 0.002 m (finest cells at shock front)",
        "h_max = 0.015 m (coarsest cells in freestream)",
        "Scaling factor alpha = 100.0",
        "Gmsh CharacteristicLengthMin = 0.002 (AMR geo)",
        "Gmsh CharacteristicLengthMax = 0.015 (AMR geo)",
        "3 AMR iterations (configurable in amr_pipeline.input)",
    ])

# ========================================================================
# SLIDE 16: Key Findings
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Key Findings")
add_footer(slide, 16, TOTAL_SLIDES)

add_info_box(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(2.5),
    "Solver Performance",
    [
        "rhoCentralFoam successfully handles Mach 3 on unstructured Gmsh mesh",
        "sonicFoam crashes at ANY supersonic speed on this mesh type",
        "Minmod limiter is essential; other limiters cause oscillations at Mach 3",
        "maxCo = 0.3 required for stability with strong bow shock",
    ])

add_info_box(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(2.5),
    "AMR Effectiveness",
    [
        "Pressure gradient captures bow shock, stagnation, and wake accurately",
        "3 AMR loops sufficient for mesh convergence (17.5k vs 4.8k nodes)",
        "Self-correcting: AMR 1 over-refines, AMR 2-3 converge to optimal",
        "Fully automated via all_run.sh with single config file",
    ])

add_info_box(slide, Inches(0.5), Inches(4.4), Inches(6), Inches(2.5),
    "Limitations Discovered",
    [
        "Mach 3.5+ crashes: negative temperature from shock oscillations",
        "Mach 27+ (re-entry): needs hy2Foam (hyStrath) for thermochemical non-equilibrium",
        "hyStrath incompatible with OpenFOAM v2406",
        "Subsonic flows (Mach 0.5): zero gradients make AMR meaningless",
    ])

add_info_box(slide, Inches(6.8), Inches(4.4), Inches(6), Inches(2.5),
    "Pipeline Features",
    [
        "Universal: works with both OpenFOAM (CFD) and COMET (DSMC)",
        "Two extraction modes: gradient (OpenFOAM) and direct (COMET MFP)",
        "Solver-agnostic: reads VTU files, not solver-specific formats",
        "Open-source ready: users configure ONLY amr_pipeline.input",
    ])

# ========================================================================
# SLIDE 17: Conclusions & Future Work
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Conclusions & Future Work")
add_footer(slide, 17, TOTAL_SLIDES)

add_info_box(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(3.5),
    "Conclusions",
    [
        "Successfully demonstrated gradient-based AMR for Mach 3 supersonic cylinder",
        "Pressure gradient is an effective adaptation criterion for shock-dominated flows",
        "The automated pipeline (Gmsh + OpenFOAM + ParaView) reduces manual meshing effort",
        "AMR achieves better shock resolution than uniform mesh with moderate cell count increase",
        "The pipeline is configurable, solver-agnostic, and ready for other flow problems",
    ])

add_info_box(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(3.5),
    "Future Work",
    [
        "Apply pipeline to 3D DSMC (COMET) results using direct MFP extraction",
        "Test with friend's RAMC-II (Mach 27) DSMC results when available",
        "Investigate multi-criteria adaptation (combine grad(p) + grad(T))",
        "Extend to higher Mach numbers with structured/hybrid meshes",
        "Integrate pipeline into thesis as AMR chapter demonstration",
    ])

# Thank you box
ty = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(2), Inches(5.5), Inches(9), Inches(1.3))
ty.fill.solid()
ty.fill.fore_color.rgb = DARK_BLUE
ty.line.fill.background()
tf = ty.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Fahim Shahriyar | fahim.shahriyar@und.edu | Advisor: Dr. Chonglin Zhang"
p2.font.size = Pt(14)
p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
p2.alignment = PP_ALIGN.CENTER

# ========================================================================
# SAVE
# ========================================================================
outpath = os.path.join(BASE, "Mach3_Cylinder_AMR_Presentation.pptx")
prs.save(outpath)
print(f"Presentation saved: {outpath}")
print(f"Total slides: {len(prs.slides)}")
